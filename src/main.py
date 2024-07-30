import os
from pptx import Presentation
from transformers import pipeline
import speech_recognition as sr

# Initialize the AI models for text generation and sentiment analysis
generator = pipeline('text-generation', model='gpt-2')
sentiment_analyzer = pipeline('sentiment-analysis')

def create_presentation(title, content_prompts, template=None):
    # Create a new PowerPoint presentation
    prs = Presentation()

    # Add a title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Generate content for each prompt and add slides
    for prompt in content_prompts:
        generated_text = generator(prompt, max_length=100, num_return_sequences=1)[0]['generated_text']
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        title_placeholder.text = prompt
        content_placeholder.text = generated_text

    # Save the presentation
    prs.save('generated_presentation.pptx')
    print("Presentation created successfully!")

def recognize_speech():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("Please say a command...")
        audio = recognizer.listen(source)
        try:
            text = recognizer.recognize_google(audio)
            print(f"You said: {text}")
            return text
        except sr.UnknownValueError:
            print("Google Speech Recognition could not understand audio")
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
    return None

def process_voice_command(command):
    if "create presentation" in command.lower():
        title = "Voice-Created Presentation"
        content_prompts = ["Introduction", "Main Points", "Conclusion"]
        create_presentation(title, content_prompts)
    elif "add slide" in command.lower():
        print("Adding a new slide...")
        # Here you would add logic to append a new slide to an existing presentation
    elif "generate content" in command.lower():
        print("Generating content...")
        # Here you would add logic to generate content for a specific slide
    else:
        print("Command not recognized. Please try again.")

def suggest_improvements(content):
    suggestions = []
    for text in content:
        sentiment = sentiment_analyzer(text)[0]
        if sentiment['label'] == 'POSITIVE':
            suggestions.append(f"The content '{text}' has a positive tone. Consider balancing it with potential challenges or limitations.")
        elif sentiment['label'] == 'NEGATIVE':
            suggestions.append(f"The content '{text}' has a negative tone. Consider highlighting some positive aspects or solutions.")
        else:
            suggestions.append(f"The content '{text}' is neutral. Consider adding more impactful or emotional language to engage the audience.")
    return suggestions

if __name__ == "__main__":
    while True:
        print("\nWhat would you like to do? (speak a command or type 'exit' to quit)")
        user_input = input().lower()

        if user_input == 'exit':
            break
        elif user_input == 'speak':
            speech_text = recognize_speech()
            if speech_text:
                process_voice_command(speech_text)
        else:
            process_voice_command(user_input)

    # Example usage of Copilot-like suggestions
    content_prompts = [
        "AI is revolutionizing industries worldwide",
        "The ethical implications of AI are concerning",
        "AI has the potential to solve complex problems"
    ]
    suggestions = suggest_improvements(content_prompts)
    print("\nHere are some suggestions for improvement:")
    for suggestion in suggestions:
        print(suggestion)
